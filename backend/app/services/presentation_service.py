"""Presentation conversion and narration section helpers."""

import math
from pathlib import Path
from typing import Dict, List

from PIL import Image, ImageDraw, ImageFont

from app.services.storage_service import STORAGE_DIR, storage_url


def _placeholder_slide(path: Path, title: str, subtitle: str) -> None:
    """Create a clean placeholder slide image when full conversion is unavailable."""
    image = Image.new("RGB", (1600, 900), "#f8fafc")
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 1600, 110), fill="#1d4ed8")
    draw.text((80, 36), "EduAvatar Presenter Studio", fill="white")
    draw.text((80, 220), title, fill="#0f172a")
    draw.text((80, 300), subtitle, fill="#475569")
    draw.rectangle((80, 680, 1520, 760), outline="#cbd5e1", width=3)
    draw.text((110, 705), "Placeholder slide image - replace this converter with production export later.", fill="#64748b")
    image.save(path)


def _convert_pdf(project_id: str, presentation_path: Path, slide_dir: Path) -> List[Path]:
    """Convert PDF pages to images using pdf2image, with a placeholder fallback."""
    try:
        from pdf2image import convert_from_path

        pages = convert_from_path(str(presentation_path), dpi=160)
        slide_paths = []
        for index, page in enumerate(pages, start=1):
            output = slide_dir / f"slide_{index:03d}.png"
            page.save(output, "PNG")
            slide_paths.append(output)
        return slide_paths
    except Exception:
        output = slide_dir / "slide_001.png"
        _placeholder_slide(output, "PDF presentation", "Install Poppler for full PDF page conversion.")
        return [output]


def _convert_pptx(project_id: str, presentation_path: Path, slide_dir: Path) -> List[Path]:
    """Create slide images from PPTX metadata; production systems can swap in LibreOffice rendering."""
    try:
        from pptx import Presentation

        deck = Presentation(str(presentation_path))
        slide_paths = []
        for index, slide in enumerate(deck.slides, start=1):
            output = slide_dir / f"slide_{index:03d}.png"
            title = f"PPTX slide {index}"
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
            subtitle = " | ".join(texts[:2]) or "Text extraction placeholder for this slide."
            _placeholder_slide(output, title, subtitle[:140])
            slide_paths.append(output)
        return slide_paths or [_fallback_presentation_slide(slide_dir)]
    except Exception:
        return [_fallback_presentation_slide(slide_dir)]


def _fallback_presentation_slide(slide_dir: Path) -> Path:
    """Create one fallback slide if conversion fails."""
    output = slide_dir / "slide_001.png"
    _placeholder_slide(output, "Presentation", "A placeholder slide was generated for preview.")
    return output


def process_presentation(project_id: str, presentation_path: Path) -> List[Dict[str, str]]:
    """Convert a PDF or PPTX into slide image records saved under slides/{project_id}/."""
    slide_dir = STORAGE_DIR / "slides" / project_id
    slide_dir.mkdir(parents=True, exist_ok=True)
    suffix = presentation_path.suffix.lower()
    if suffix == ".pdf":
        slide_paths = _convert_pdf(project_id, presentation_path, slide_dir)
    elif suffix == ".pptx":
        slide_paths = _convert_pptx(project_id, presentation_path, slide_dir)
    else:
        slide_paths = [_fallback_presentation_slide(slide_dir)]
    return [
        {"index": index, "path": str(path), "url": storage_url(path)}
        for index, path in enumerate(slide_paths, start=1)
    ]


def build_script_sections(
    project_id: str,
    original_script: str,
    target_language: str,
    slide_count: int,
) -> List[Dict[str, object]]:
    """Divide script text into slide-sized narration chunks and estimate duration."""
    words = original_script.strip().split()
    chunk_size = max(1, math.ceil(len(words) / slide_count))
    sections = []
    for index in range(slide_count):
        chunk_words = words[index * chunk_size : (index + 1) * chunk_size]
        if not chunk_words:
            chunk_words = [f"Slide {index + 1} narration placeholder."]
        text = " ".join(chunk_words)
        estimated_duration = max(4.0, round(len(text.split()) / 2.4, 1))
        sections.append(
            {
                "slide_index": index + 1,
                "text": text,
                "target_language": target_language,
                "estimated_duration_seconds": estimated_duration,
            }
        )
    return sections
